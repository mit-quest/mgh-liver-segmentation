import cv2
import numpy as np
import matplotlib.pyplot as plt
import math
import csv
import pandas as pd
import argparse
from typing import List, Dict, Any
from statistics import median

from skimage import io, data, color, img_as_ubyte
from skimage.color import rgb2gray
from skimage.morphology import binary_opening, disk, remove_small_holes, remove_small_objects, binary_erosion, binary_dilation
from skimage.transform import rescale, resize
from skimage.feature import canny
from skimage.draw import ellipse_perimeter, line
from skimage.measure import find_contours
from skimage.util import img_as_ubyte
from skimage import color, img_as_float
import queue
import copy
from pptx import Presentation
from pptx.util import Inches, Pt

from scipy import ndimage as ndi
from scipy.spatial import distance
from PIL import Image

import os
import sys
sys.setrecursionlimit(10000000)

import imagej
import cv2

SHARPEN_FILTER = np.array([[0, -1, 0],
                           [-1, 5, -1],
                           [0, -1, 0]])

# Slide title dimensions
title_left = Inches(4.4)
title_top = Inches(0)
title_width = Inches(1.2)
title_height = Inches(0.65)
title_font_size = Pt(36)

# Slide image dimensions
image_left1 = Inches(0.5)
image_top1 = Inches(0.65)
image_left2 = Inches(5.3)
image_top2 = Inches(3.8)
image_width = Inches(2)
image_horizontal_margin = Inches(0.2)

# Slide fat estimate dimensions
fat_estimate_left1 = Inches(3.35)
fat_estimate_top1 = Inches(3.3)
fat_estimate_left2 = Inches(8.1)
fat_estimate_top2 = Inches(6.45)
fat_estimate_width = Inches(0.7)
fat_estimate_height = Inches(0.4)
fat_estimate_font_size = Pt(20)

# Slide bottom text dimensions
bottom_text_top = Inches(6.95)
bottom_text_left = Inches(0.4)
bottom_text_width = Inches(3.3)
bottom_text_height = Inches(0.5)

# Map index id to slide dimensions
image_left_map = {0: image_left1, 1: image_left2, 2: image_left1, 3: image_left2}
image_top_map = {0: image_top1, 1: image_top1, 2: image_top2, 3: image_top2}
fat_estimate_left_map = {0: fat_estimate_left1, 1: fat_estimate_left2, 2: fat_estimate_left1, 3: fat_estimate_left2}
fat_estimate_top_map = {0: fat_estimate_top1, 1: fat_estimate_top1, 2: fat_estimate_top2, 3: fat_estimate_top2}


class Graph:
    def __init__(self, row, col, graph):
        self.ROW = row
        self.COL = col
        self.graph = graph

    def isSafe(self, i, j, visited):
        # Row number is in range, column number is in range, value is 1 and not yet visited
        return 0 <= i and i < self.ROW and 0 <= j and j < self.COL and not visited[i][j] and self.graph[i][j]

    def BFS(self, i, j, visited, island):
        # Utility function to do BFS for a 2D boolean matrix. Uses only the 4 neighbors as adjacent vertices
        rowNbr = [-1, 0, 1, 0]
        colNbr = [0, -1, 0, 1]
        q = []
        q.append((i,j))
        visited[i][j] = True

        while len(q) != 0:
            x,y = q.pop(0)
            for k in range(len(rowNbr)):
                if self.isSafe(x + rowNbr[k], y + colNbr[k], visited):
                    island.append((x + rowNbr[k], y + colNbr[k]))
                    visited[(x) + rowNbr[k]][y + colNbr[k]] = True
                    q.append((x + rowNbr[k], y + colNbr[k]))

    def findIslands(self):
        # Make a bool array to mark visited cells. Initially all cells are unvisited
        visited = [[False for j in range(self.COL)]for i in range(self.ROW)]
        # Initialize count as 0 and traverse through cells of given matrix
        index = 0
        islands = []
        for i in range(self.ROW):
            for j in range(self.COL):
                # If a cell with value 1 is not visited yet, then new island found
                if visited[i][j] == False and self.graph[i][j] == 1:
                    # Visit all cells in this island and increment island count
                    island = []
                    self.BFS(i, j, visited, island)
                    islands.append(island)
                    index += 1
        return islands


def sharpen_image(image_path, filter, show=False):
    original_image = cv2.imread(image_path)
    sharpened_image = cv2.filter2D(original_image, -1, filter)
    return original_image, sharpened_image


def remove_background(opened_image_bool):
    np_graph = opened_image_bool
    row, col = np_graph.shape
    graph = np_graph.tolist()
    g = Graph(row, col, graph)
    islands = g.findIslands()

    # Filter islands by size. Keep islands with sizes less than max_size
    island_to_score_map = dict()
    for island in islands:
        if len(island) < 1000: # 1000 good for formalin HF-3 liver
            island_to_score_map[tuple(island)] = 1 # Score of 1 means keep
        else:
            island_to_score_map[tuple(island)] = 0 # Score of 0 means do not keep

    # Create binary image without obvious non-fat. Color black non-fat islands
    binary_without_obvious_non_fat = np.zeros(opened_image_bool.shape)
    for island, score in island_to_score_map.items():
        if score == 1:
            for x, y in island:
                binary_without_obvious_non_fat[x][y] = 255
    binary_without_obvious_non_fat = binary_without_obvious_non_fat.astype(np.uint8)
    return binary_without_obvious_non_fat


# Assign number between 0 (most likely non-fat) and 1 (most likely fat) for
# probability that island is fat. Input is list of lists of pixels for each island
def get_fat_score_for_watershed_image(islands, original_mask, min_size=10, max_size=1000):
    # Filter islands. Map of tuple of (x, y) coordinates of island to
    # probability that island is fat, initialized to 0.5 (undecided)
    island_to_score_map = dict()
    for island in islands:
        if min_size <= len(island) < max_size: # Filter out small and large tears and holes
            island_to_score_map[tuple(island)] = 0.5
        else:
            island_to_score_map[tuple(island)] = 0 # Small and large tears and holes are assumed to be non-fat (0)

    # --- Assign circularity to islands ---
    island_to_circularity_map = dict() # Map of tuple of (x, y) coordinates of island to circularity score
    for island in island_to_score_map:
        if island_to_score_map[island] > 0:
            # Create mask with 1 for islands and 0 for non-island
            island_mask = np.zeros(original_mask.shape)
            for x, y in island:
                island_mask[x][y] = 1
            island_mask = island_mask.astype(np.uint8)
            contours, hierarchy = cv2.findContours(island_mask, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_NONE)

            # Find circularity of island
            contour = contours[0]
            contour_area = cv2.contourArea(contour)
            contour_perimeter = cv2.arcLength(contour, True)
            center, radius = cv2.minEnclosingCircle(contour)

            # --- Method to calculate circularity: contour area vs. contour perimeter ---
            circularity = (2 * math.pi * math.sqrt(contour_area / math.pi)) / contour_perimeter

            # Assign circularity
            island_to_circularity_map[island] = circularity

    # --- Filter islands by circularity ---
    for island, circularity in island_to_circularity_map.items():
        if circularity >= 0.67:
            island_to_score_map[island] = 1
        if circularity < 0.67:
            island_to_score_map[island] = 0

    # Create new mask
    new_mask = np.zeros(original_mask.shape)
    for island, circularity in island_to_circularity_map.items():
        if island_to_score_map[island] == 1: # 0 means non-fat, 0.5 means unsure, 1 means fat
            for x, y in island:
                new_mask[x][y] = 255

    new_mask = new_mask.astype(np.uint8)
    return new_mask


def count_fat_macro_micro(image, mask):
    image_255 = copy.deepcopy(image)*255 # Use cv2_imshow to show image
    image_inv = copy.deepcopy(image)

    # Find area of black background in original image
    num_black_border_pixels = 0
    for row in range(len(image_255)):
        for col in range(len(image_255[row])):
            if image_255[row][col] <= 45 and not (170 <= row <= 520 and 70 <= col <= 410): # 45 experimentally determined
                image_inv[row][col] = 255
            else:
                image_inv[row][col] = 0
    _, binary_image_inv = cv2.threshold(image_inv, 200, 255, cv2.THRESH_BINARY) # 200 was chosen to separate 0 (black) and 255 (white) pixels
    binary_image_inv = binary_opening(binary_image_inv, disk(1))
    image_np_graph_inv = binary_image_inv
    image_row_inv, image_col_inv = image_np_graph_inv.shape
    image_graph_inv = image_np_graph_inv.tolist()
    image_g_inv = Graph(image_row_inv, image_col_inv, image_graph_inv)
    image_islands_inv = image_g_inv.findIslands()
    for island in image_islands_inv:
        if len(island) > 1000: # 1000 was chosen as an arbitrarily large number of pixels
            num_black_border_pixels += len(island)

    # Find area of large white islands
    num_white_island_pixels = 0
    _, binary_image = cv2.threshold(image_255, 200, 255, cv2.THRESH_BINARY) # 150 experimentally determined
    binary_image = binary_opening(binary_image, disk(1))
    image_np_graph = binary_image
    image_row, image_col = image_np_graph.shape
    image_graph = image_np_graph.tolist()
    image_g = Graph(image_row, image_col, image_graph)
    image_islands = image_g.findIslands()
    for island in image_islands:
        if len(island) > 700: # 700 was experimentally determined
            num_white_island_pixels += len(island)

    # Find liver area excluding border and large tears
    liver_area = (image_255.shape[0] * image_255.shape[1]) - num_black_border_pixels - num_white_island_pixels

    # Find fat area using predicted mask
    fat_areas = []
    mask_255 = copy.deepcopy(mask)*255
    mask_macro = copy.deepcopy(mask) # Not needed but helpful for debugging
    _, binary_mask = cv2.threshold(mask_255, 30, 255, cv2.THRESH_BINARY) # 30 experimentally determined
    binary_mask = binary_opening(binary_mask, disk(1))
    mask_np_graph = binary_mask
    mask_row, mask_col = mask_np_graph.shape
    mask_graph = mask_np_graph.tolist()
    mask_g = Graph(mask_row, mask_col, mask_graph)
    mask_islands = mask_g.findIslands()
    for island in mask_islands:
        fat_areas.append(len(island))
        if len(island) >= 30:
            for row, col in island:
                mask_macro[row][col] = 255
        else:
            for row, col in island:
                mask_macro[row][col] = 150
    num_fat = len(fat_areas)
    total_fat_area = sum(fat_areas)

    macro = []
    micro = []
    for fat_area in fat_areas:
        if fat_area >= 30:
            macro.append(fat_area)
        else:
            micro.append(fat_area)
    num_macro = len(macro)
    num_micro = len(micro)
    total_macro_area = sum(macro)
    total_micro_area = sum(micro)

    # Calculate fat estimates with total liver area
    total_fat_percentage = total_fat_area / liver_area * 100
    macro_fat_percentage = total_macro_area / liver_area * 100
    return total_fat_percentage, macro_fat_percentage


def mean_fat_percent(image_names, original_folder, new_mask_folder):
    # Get original images
    image_list = []
    for liver_file in image_names:
        image_file = os.path.join(original_folder, liver_file)
        original_image = Image.open(image_file).convert('L')
        if original_image.size != (480, 640):
            image_list.append(np.array(original_image.resize((480,640))))
        else:
            image_list.append(np.array(original_image))
    image_np = np.asarray(image_list)
    images = np.asarray(image_np, dtype=np.float32)/255
    images = images.reshape(images.shape[0], images.shape[1], images.shape[2], 1)

    # Get new masks
    mask_list = []
    for liver_file in image_names:
        mask_file = os.path.join(new_mask_folder, liver_file)
        new_mask = Image.open(mask_file).convert('L')
        if new_mask.size != (480, 640):
            mask_list.append(np.array(new_mask.resize((480,640))))
        else:
            mask_list.append(np.array(new_mask))
    mask_np = np.asarray(mask_list)
    masks = np.asarray(mask_np, dtype=np.float32)/255
    masks = masks.reshape(masks.shape[0], masks.shape[1], masks.shape[2], 1)

    # Find mean total and macro fat percent
    liver_total_fat = []
    liver_macro_fat = []
    num_liver_files = len(image_names)
    for i in range(num_liver_files):
        total_fat, macro_fat = count_fat_macro_micro(images[i], masks[i])
        liver_total_fat.append(total_fat)
        liver_macro_fat.append(macro_fat)

    # Calculate min, max, mean, and median total fat
    min_total_fat = min(liver_total_fat)
    max_total_fat = max(liver_total_fat)
    mean_total_fat = sum(liver_total_fat) / num_liver_files
    median_total_fat = median(liver_total_fat)

    # Calculate min, max, mean, and median macro fat
    min_macro_fat = min(liver_macro_fat)
    max_macro_fat = max(liver_macro_fat)
    mean_macro_fat = sum(liver_macro_fat) / num_liver_files
    median_macro_fat = median(liver_macro_fat)

    return round(mean_total_fat, 2), round(mean_macro_fat, 2)


def run_watershed(binary_image_path, watershed_image_path, ij):
    macro = """
    #@ String binary_image_path
    #@ String watershed_image_path

    open(binary_image_path);
    run("8-bit");
    run("Watershed");
    saveAs("Tiff", watershed_image_path);
    """
    args = {
        'binary_image_path': binary_image_path,
        'watershed_image_path': watershed_image_path,
    };
    result = ij.py.run_macro(macro, args);
    result.getOutput('watershed_image_path')


def liver_segmentation(images_directory, output_directory, liver_name, image_name, ij):
    image_path = os.path.join(images_directory, liver_name, image_name)

    # Find sharpened, grayscale, and binary images
    original_image, sharpened_image = sharpen_image(image_path, SHARPEN_FILTER, show=True)
    sharpened_gray = cv2.cvtColor(sharpened_image, cv2.COLOR_BGR2GRAY)
    _, sharpened_binary = cv2.threshold(sharpened_gray, 205, 255, cv2.THRESH_BINARY) # 205 good for formalin HF-3

    # Remove noise
    opened_image_bool = remove_small_holes(sharpened_binary, area_threshold=5) # Change black area < 5 pixels to white
    opened_image_bool = remove_small_objects(opened_image_bool, min_size=10) # Change white area < 10 pixels to black
    opened_image = opened_image_bool.astype(np.uint8)  # Convert to an unsigned byte
    opened_image *= 255

    # Find binary image without background, save binary
    binary_without_obvious_non_fat = remove_background(opened_image_bool)
    save_path = os.path.join(output_directory, liver_name, image_name)
    cv2.imwrite(save_path, binary_without_obvious_non_fat)

    # Run watershed and save image
    run_watershed(save_path, save_path, ij)

    # After applying watershed
    watershed_binary = cv2.imread(save_path)
    watershed_gray = cv2.cvtColor(watershed_binary, cv2.COLOR_BGR2GRAY)
    watershed_bool = watershed_gray > 0

    # Apply erosion
    erode_image_bool = binary_erosion(watershed_gray)
    erode_image = erode_image_bool.astype(np.uint8)
    erode_image *= 255

    # Create mask
    np_graph = erode_image_bool # Use eroded watershed binary
    row, col = np_graph.shape
    graph = np_graph.tolist()
    g = Graph(row, col, graph)
    islands = g.findIslands()

    # Create new mask, convert to overlay with green mask
    new_mask = get_fat_score_for_watershed_image(islands, np_graph, min_size=2, max_size=1000) # With erosion, adjust size thresholds
    new_mask_dilate = binary_dilation(new_mask)
    new_mask_float = img_as_float(new_mask_dilate)
    green_multiplier = [0, 1, 0]
    new_mask_rgb = color.gray2rgb(new_mask_float) * green_multiplier

    # Save mask
    new_mask_to_save = new_mask_rgb.astype(np.uint8) * 255
    cv2.imwrite(save_path, new_mask_to_save)

    # Estimate fat
    original_image_path = os.path.join(images_directory, liver_name)
    mask_image_path = os.path.join(output_directory, liver_name)
    mean_total_fat, mean_macro_fat = mean_fat_percent([image_name], original_image_path, mask_image_path)

    # Save fat estimates per liver to CSV file. Save in order of image name, total fat, macro fat (arbitrary threshold)
    csv_save_path = os.path.join(output_directory, liver_name, f'{liver_name}_fat_estimates.csv')
    with open(csv_save_path, 'a') as f:
        writer = csv.writer(f)
        writer.writerow([image_name, mean_total_fat, mean_macro_fat])

    background_path = os.path.join(original_image_path, image_name)
    background = Image.open(background_path)
    overlay_path = os.path.join(output_directory, liver_name, image_name)
    overlay = Image.open(overlay_path)
    background = background.convert("RGBA")
    overlay = overlay.convert("RGBA")
    new_img = Image.blend(background, overlay, 0.5)
    new_img.save(overlay_path, "tiff")


def make_powerpoint(images_directory, output_directory, pathologist_estimates, liver_name, powerpoint_save_path):
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    liver_folder_path = os.path.join(output_directory, liver_name)
    liver_files = os.listdir(liver_folder_path)
    fat_estimates_path = os.path.join(liver_folder_path, f'{liver_name}_fat_estimates.csv')

    # Find overall fat estimate for liver
    df = pd.read_csv(fat_estimates_path, header=None)
    num_rows = len(df.index)
    liver_sum_estimates = 0
    liver_num_estimates = 0
    for index, row in df.iterrows():
        total_fat = row[1]
        liver_sum_estimates += float(total_fat)
        liver_num_estimates += 1
    liver_overall_estimate = liver_sum_estimates / liver_num_estimates

    # Get pathologist estimates for liver
    pathologist_estimates_df = pd.read_csv(pathologist_estimates)
    pathologist_liver_estimates = 'N/A'
    for index, row in pathologist_estimates_df.iterrows():
        liver = row[0]
        if liver == liver_name:
            liver_fat_combined = ''
            for i in range(1, len(row)):
                if len(liver_fat_combined) == 0:
                    liver_fat_combined = row[i]
                else:
                    liver_fat_combined += ', ' + row[i]
            pathologist_liver_estimates = liver_fat_combined
            break

    # Create slides
    for index, row in df.iterrows():
        image_name = row[0]
        total_fat = row[1]
        macro_fat = row[2]
        original_image_path = os.path.join(images_directory, liver_name, image_name)
        mask_path = os.path.join(output_directory, liver_name, image_name)
        index_id = index % 4

        if index % 4 == 0:
            # Add liver name as title
            slide = prs.slides.add_slide(blank_slide_layout)
            shapes = slide.shapes
            txBox = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = liver_name
            p.font.size = title_font_size

        pic = slide.shapes.add_picture(original_image_path, image_left_map[index_id], image_top_map[index_id], width=image_width)
        pic = slide.shapes.add_picture(mask_path, image_left_map[index_id]+image_width+image_horizontal_margin, image_top_map[index_id], width=image_width)
        txBox = slide.shapes.add_textbox(fat_estimate_left_map[index_id], fat_estimate_top_map[index_id], fat_estimate_width, fat_estimate_height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f'{total_fat}%'
        p.font.size = fat_estimate_font_size

        if index % 4 == 3 or index == num_rows-1:
            # Add bottom textbox with overall estimates
            txBox = slide.shapes.add_textbox(bottom_text_left, bottom_text_top, bottom_text_width, bottom_text_height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            fat_estimates = f'Pathologist macro estimate: {pathologist_liver_estimates}% | Algorithm total estimate: {round(liver_overall_estimate, 2)}%'
            p.text = fat_estimates
            p.font.size = Pt(20)

    prs.save(powerpoint_save_path)


def parse_args(args: List[str]) -> Dict[str, Any]:
    parser = argparse.ArgumentParser(description="Estimate fat in liver biopsy images")
    parser.add_argument(
        "--images_directory",
        type=str,
        required=True,
        help="Path to image directory",
    )
    parser.add_argument(
        "--output_directory",
        type=str,
        required=True,
        help="Path to output directory to save segmentation and fat estimate",
    )
    parser.add_argument(
        "--pathologist_estimates",
        type=str,
        required=True,
        help="Path to CSV with pathologist estimates",
    )
    parser.add_argument(
        "--magnification",
        type=str,
        required=True,
        help="Magnification of images to use, e.g. 20x",
    )
    args = vars(parser.parse_args())
    return args


def main(**args: Dict[str, Any]) -> None:
    images_directory, output_directory, pathologist_estimates, magnification = args.values()
    ij = imagej.init()
    liver_folders = os.listdir(images_directory)

    # Create output folder
    if not os.path.exists(output_directory):
        os.mkdir(output_directory)

    for liver_name in liver_folders:
        # Create liver output folder
        liver_output_folder = os.path.join(output_directory, liver_name)
        if not os.path.exists(liver_output_folder):
            os.mkdir(liver_output_folder)

        liver_folder_path = os.path.join(images_directory, liver_name)
        liver_files = os.listdir(liver_folder_path)
        liver_images = []

        csv_file_path = os.path.join(liver_output_folder, f'{liver_name}_fat_estimates.csv')
        if os.path.exists(csv_file_path):
            os.remove(csv_file_path)

        for file in liver_files:
            if file.endswith(".tiff") and magnification in file: # Could use REGEX
                liver_images.append(file)

        for image_name in liver_images:
                liver_segmentation(images_directory, output_directory, liver_name, image_name, ij)

        # Create powerpoint for each liver
        powerpoint_save_path = os.path.join(liver_output_folder, f'{liver_name}_slides.pptx')
        make_powerpoint(images_directory, output_directory, pathologist_estimates, liver_name, powerpoint_save_path)


if __name__ == "__main__":
    kwargs = parse_args(sys.argv[1:])
    main(**kwargs)
