import argparse
import csv
from typing import List, Dict, Any
import os
import pandas as pd
import scipy.stats as stats

from pptx import Presentation
from pptx.util import Inches, Pt


def parse_args(args: List[str]) -> Dict[str, Any]:
    parser = argparse.ArgumentParser(description="Estimate fat in liver biopsy images")
    parser.add_argument(
        "--images_directory",
        type=str,
        required=True,
        help="Path to image directory",
    )
    parser.add_argument(
        "--output_directory",
        type=str,
        required=True,
        help="Path to output directory to save segmentation and fat estimate",
    )
    parser.add_argument(
        "--magnification",
        type=str,
        required=True,
        help="Magnification of images to use ('10x', '20x', or '40x)",
    )
    parser.add_argument(
        "--preservation",
        type=str,
        required=True,
        help="Preservation type ('formalin' or 'frozen')",
    )
    parser.add_argument(
        "--pathologist_estimates",
        type=str,
        required=False,
        help="Path to CSV with pathologist estimates",
    )
    args = vars(parser.parse_args())
    return args

def calculate_liver_fat(fat_estimates_path):

    df = pd.read_csv(fat_estimates_path, header=None)
    num_rows = len(df.index)
    liver_sum_estimates = []
    for index, row in df.iterrows():
        total_fat = row[1]
        liver_sum_estimates.append(float(total_fat))

    # drop image estimates if > 1 SD
    zscore_arr = stats.zscore(liver_sum_estimates, axis=0)
    for i, e in enumerate(liver_sum_estimates):
        if zscore_arr[i] >= 2:
            print("dropping " + str(e))
    new_liver_sum_estimates = [e for i, e in enumerate(liver_sum_estimates) if zscore_arr[i] <= 2]

    liver_overall_mean_estimate = sum(new_liver_sum_estimates) / len(new_liver_sum_estimates)

    with open(fat_estimates_path, 'a') as f:
        writer = csv.writer(f)
        f.write("\n")
        writer.writerow(['Total liver mean fat estimate', round(liver_overall_mean_estimate, 2)])

    return liver_overall_mean_estimate

# TODO: Center title and pathologist/algorithm estimates
def make_powerpoint(images_directory, output_directory, pathologist_estimates, liver_name, powerpoint_save_path):
    
    # Slide title dimensions
    title_left = Inches(4.4)
    title_top = Inches(0)
    title_width = Inches(1.2)
    title_height = Inches(0.65)
    title_font_size = Pt(36)

    # Slide image dimensions
    image_left1 = Inches(0.5)
    image_top1 = Inches(0.65)
    image_left2 = Inches(5.3)
    image_top2 = Inches(3.8)
    image_width = Inches(2)
    image_horizontal_margin = Inches(0.2)

    # Slide fat estimate dimensions
    fat_estimate_left1 = Inches(3.35)
    fat_estimate_top1 = Inches(3.3)
    fat_estimate_left2 = Inches(8.1)
    fat_estimate_top2 = Inches(6.45)
    fat_estimate_width = Inches(0.7)
    fat_estimate_height = Inches(0.4)
    fat_estimate_font_size = Pt(20)

    # Slide bottom text dimensions
    bottom_text_top = Inches(6.95)
    bottom_text_left = Inches(0.4)
    bottom_text_width = Inches(3.3)
    bottom_text_height = Inches(0.5)

    # Map index id to slide dimensions
    image_left_map = {0: image_left1, 1: image_left2, 2: image_left1, 3: image_left2}
    image_top_map = {0: image_top1, 1: image_top1, 2: image_top2, 3: image_top2}
    fat_estimate_left_map = {0: fat_estimate_left1, 1: fat_estimate_left2, 2: fat_estimate_left1, 3: fat_estimate_left2}
    fat_estimate_top_map = {0: fat_estimate_top1, 1: fat_estimate_top1, 2: fat_estimate_top2, 3: fat_estimate_top2}

    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    liver_folder_path = os.path.join(output_directory, liver_name)
    liver_files = os.listdir(liver_folder_path)
    fat_estimates_path = os.path.join(liver_folder_path, f'{liver_name}_fat_estimates.csv')

    # Find overall fat estimate for liver & write it to csv
    liver_overall_estimate = calculate_liver_fat()

    # Get pathologist estimates for liver
    pathologist_estimates_df = pd.read_csv(pathologist_estimates)
    pathologist_liver_estimates = 'N/A'
    for index, row in pathologist_estimates_df.iterrows():
        liver = row[0]
        if liver == liver_name:
            liver_fat_combined = ''
            for i in range(1, len(row)):
                if len(liver_fat_combined) == 0:
                    liver_fat_combined = row[i]
                else:
                    liver_fat_combined += ', ' + row[i]
            pathologist_liver_estimates = liver_fat_combined
            break

    # Create slides
    for index, row in df.iterrows():
        image_name = row[0]
        total_fat = row[1]
        macro_fat = row[2]
        original_image_path = os.path.join(images_directory, liver_name, image_name)
        mask_path = os.path.join(output_directory, liver_name, image_name)
        index_id = index % 4

        if index % 4 == 0:
            # Add liver name as title
            slide = prs.slides.add_slide(blank_slide_layout)
            shapes = slide.shapes
            txBox = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = liver_name
            p.font.size = title_font_size

        pic = slide.shapes.add_picture(original_image_path, image_left_map[index_id], image_top_map[index_id], width=image_width)
        pic = slide.shapes.add_picture(mask_path, image_left_map[index_id]+image_width+image_horizontal_margin, image_top_map[index_id], width=image_width)
        txBox = slide.shapes.add_textbox(fat_estimate_left_map[index_id], fat_estimate_top_map[index_id], fat_estimate_width, fat_estimate_height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f'{total_fat}%'
        p.font.size = fat_estimate_font_size

        if index % 4 == 3 or index == num_rows-1:
            # Add bottom textbox with overall estimates
            txBox = slide.shapes.add_textbox(bottom_text_left, bottom_text_top,
                                             bottom_text_width, bottom_text_height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            fat_estimates = f'Pathologist macro estimate: {pathologist_liver_estimates}% | Algorithm total estimate: {round(liver_overall_estimate, 2)}%'
            p.text = fat_estimates
            p.font.size = Pt(20)

    prs.save(powerpoint_save_path)